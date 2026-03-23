# scripts/rebuild_deck.py
"""
Rebuilds the IOCL pitch deck from scratch.
Run: /Users/abindal/opt/anaconda3/bin/python3 scripts/rebuild_deck.py
Output: asset/Santosh_IOCL_Pitch_Deck_v2.pptx
"""

import sys
sys.path.insert(0, ".")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from scripts.deck_helpers import (
    C_BG, C_TEAL, C_TEAL2, C_WHITE, C_MUTED, C_CARD, C_CARD_BORDER,
    C_FOOTER, C_GREY, W, H, IMG_W, CONTENT_X, CONTENT_W, FOOTER_Y, FONT,
    add_bg, add_rect, add_text, add_text_multiline, add_image, add_footer,
    add_logo, add_slide_tag, add_headline, add_subheadline, add_kpi_chip,
    add_bullet_row, add_card, add_pattern_a_image,
    download_image, get_logo_bytes
)


def new_slide(prs):
    """Add a blank slide to the presentation."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER (Pattern A)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_01(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_pattern_a_image(slide, images.get("photo-1565008576549-57569a49371d"))

    # Logo + company name (top of right panel)
    add_logo(slide, logo)
    add_text(slide, "SANTOSH PETROCHEMICAL INNOVATIONS PVT. LTD.",
        CONTENT_X, 0.18, CONTENT_W - 0.7, 0.32,
        size=7.5, color=C_MUTED, bold=False)

    # Tag
    add_text(slide, "PARTNERSHIP PROPOSAL  ·  MARCH 2026",
        CONTENT_X, 0.62, CONTENT_W, 0.22,
        size=8, bold=True, color=C_TEAL)

    # Headline
    add_text(slide, "Group II+\nRe-Refined Base Oil",
        CONTENT_X, 0.88, CONTENT_W, 1.5,
        size=30, bold=True, color=C_WHITE)

    # Subheadline
    add_text(slide,
        "Presented to Indian Oil Corporation Limited\n65 TPD Plant  ·  Ghaziabad, UP",
        CONTENT_X, 2.46, CONTENT_W, 0.65,
        size=10.5, color=C_MUTED)

    # 4 KPI chips (2×2 grid)
    chips = [
        ("₹159 Cr", "Total Investment"),
        ("65 TPD", "Plant Capacity"),
        ("35 Yrs", "IOCL Relationship"),
        ("₹169 Cr+", "Revenue by Year 4"),
    ]
    chip_w, chip_h = 3.3, 0.82
    gap = 0.12
    cx = CONTENT_X
    for i, (val, lbl) in enumerate(chips):
        col = i % 2
        row = i // 2
        x = cx + col * (chip_w + gap)
        y = 3.3 + row * (chip_h + gap)
        add_kpi_chip(slide, val, lbl, x, y, width=chip_w, height=chip_h)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — REFINING PROCESS (Pattern B — Full-Width)
# ═══════════════════════════════════════════════════════════════════════════════

STEPS = [
    ("01", "Collection &\nPre-Screening",
     "Used oil collected. Feed filtered to 150–250 micron at tank farm.", False),
    ("02", "Dehydration",
     "Heated to 120°C. Water, glycol, solvents removed under vacuum.", False),
    ("03", "Vacuum\nDistillation",
     "Heated 280–300°C. Two-step evaporation separates LOBS.", True),
    ("04", "Hydrotreating",
     "H₂ + metallic catalysts at high pressure. Sulfur → <300 ppm.", True),
    ("05", "Fractionation",
     "Cuts separated: SN 150, SN 500, Bright Stock equivalent.", False),
    ("06", "QC & Dispatch",
     "Viscosity, sulfur, VI, flash point tested. Bulk / drum / IBC.", False),
]

def build_slide_07(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_logo(slide, logo)

    # Tag + headline
    add_text(slide, "TECHNOLOGY & PROCESS",
        0.35, 0.22, 10.0, 0.25, size=8, bold=True, color=C_TEAL)
    add_text(slide, "6-Step Distillation & Hydrotreating",
        0.35, 0.50, 9.5, 0.55, size=22, bold=True, color=C_WHITE)
    add_text(slide, "REVA Process Technologies — European Technology, Indian Scale",
        0.35, 1.05, 9.5, 0.28, size=10, color=C_MUTED)

    # 6 step boxes
    box_w = (W - 0.7) / 6 - 0.06
    box_h = 3.5
    box_y = 1.42
    box_x0 = 0.35

    for i, (num, name, desc, highlight) in enumerate(STEPS):
        bx = box_x0 + i * (box_w + 0.06)
        fill = RGBColor(0x09, 0x2e, 0x22) if highlight else C_CARD
        border = C_TEAL if highlight else C_CARD_BORDER
        # Top accent bar
        add_rect(slide, bx, box_y, box_w, 0.06,
                 fill=C_TEAL if highlight else RGBColor(0x00, 0x3d, 0x2e))
        add_rect(slide, bx, box_y + 0.06, box_w, box_h - 0.06,
                 fill=fill, border=border)
        add_text(slide, num,
            bx + 0.1, box_y + 0.15, box_w - 0.2, 0.32,
            size=13, bold=True,
            color=C_TEAL if highlight else RGBColor(0x00, 0x8a, 0x6e))
        add_text(slide, name,
            bx + 0.1, box_y + 0.48, box_w - 0.2, 0.65,
            size=9.5, bold=True, color=C_WHITE)
        add_text(slide, desc,
            bx + 0.1, box_y + 1.18, box_w - 0.2, 1.8,
            size=8.5, color=C_MUTED, wrap=True)

        # Arrow connector (not last)
        if i < 5:
            arrow_x = bx + box_w + 0.005
            add_text(slide, "▶",
                arrow_x, box_y + box_h/2 - 0.12, 0.06, 0.25,
                size=8, color=C_TEAL, align=PP_ALIGN.CENTER)

    # Output strip
    strip_y = box_y + box_h + 0.18
    add_rect(slide, 0.35, strip_y, W - 0.7, 0.55,
             fill=RGBColor(0x04, 0x1d, 0x13), border=C_CARD_BORDER)
    add_text(slide, "OUTPUT →",
        0.48, strip_y + 0.12, 0.9, 0.28,
        size=8, color=C_MUTED, bold=True)
    outputs = [
        "API Group II+ RRBO",
        "VI ≥ 95",
        "Sulfur <300 ppm",
        "Flash Point >200°C",
    ]
    for i, out in enumerate(outputs):
        bg = C_TEAL if i == 0 else C_CARD
        fg = C_BG if i == 0 else C_MUTED
        pill_w = 1.8 if i == 0 else 1.45
        px = 1.42 + i * 1.90
        add_rect(slide, px, strip_y + 0.09, pill_w, 0.35,
                 fill=bg, border=None)
        add_text(slide, out,
            px + 0.06, strip_y + 0.12, pill_w - 0.12, 0.28,
            size=8.5, bold=(i == 0), color=fg, align=PP_ALIGN.CENTER)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ABOUT US / GROUP HISTORY (Pattern A)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_02(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_pattern_a_image(slide, images.get("photo-1504328345606-18bbc8c9d7d1"))
    add_logo(slide, logo)
    add_slide_tag(slide, "GROUP HISTORY")
    add_headline(slide, "Two Decades of Partnership\nwith India's Oil Majors", top=0.55, size=22)

    cards = [
        ("IOCL SSI Stockist  (Est. 2003)",
         ["SERVO lubricants distributor across 6 districts of Western UP",
          "One of IOCL's longest-standing petroleum product partners"]),
        ("HPCL LPG Bottling  (Since Dec 2018)",
         ["Private bottler at Amroha, UP — 30 TMT/yr cylinder bottling",
          "Distributing for HPCL dealers across Amroha region"]),
        ("HPCL O&M Contract  (Sitarganj Plant)",
         ["5-year O&M contract for HPCL Sitarganj LPG Plant",
          "340 TMT total bottling across contract tenure"]),
    ]
    card_h = 1.35
    card_y0 = 1.85
    for i, (title, lines) in enumerate(cards):
        add_card(slide, CONTENT_X, card_y0 + i * (card_h + 0.12),
                 CONTENT_W, card_h, title, lines,
                 title_size=10, line_size=9)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ABOUT THE COMPANY (Pattern A)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_03(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_pattern_a_image(slide, images.get("photo-1518709268805-4e9042af9f23"))
    add_logo(slide, logo)
    add_slide_tag(slide, "COMPANY OVERVIEW")
    add_headline(slide, "India's Next-Gen\nGroup II+ RRBO Producer", top=0.55, size=22)

    facts = [
        ("Company", "Santosh Petrochemical Innovations Pvt. Ltd."),
        ("Technology", "REVA Process Technologies, Pune (European)"),
        ("Plant Location", "Western Uttar Pradesh (Ghaziabad)"),
        ("Capacity", "65 TPD feed (designed for 200 TPD expansion)"),
        ("Output", "API Group II+ Re-Refined Base Oil (RRBO/LOBS)"),
        ("By-products", "Fuel oil, Asphalt Extender (Residue)"),
    ]
    row_h = 0.42
    row_y0 = 1.70
    for i, (label, value) in enumerate(facts):
        y = row_y0 + i * row_h
        # Subtle row separator
        if i % 2 == 0:
            add_rect(slide, CONTENT_X, y, CONTENT_W, row_h,
                     fill=RGBColor(0x12, 0x1e, 0x2c))
        add_text(slide, label,
            CONTENT_X + 0.1, y + 0.08, 1.4, row_h - 0.1,
            size=8.5, bold=True, color=C_TEAL)
        add_text(slide, value,
            CONTENT_X + 1.55, y + 0.08, CONTENT_W - 1.65, row_h - 0.1,
            size=8.5, color=C_WHITE)

    # EPR callout
    epr_y = row_y0 + len(facts) * row_h + 0.18
    add_rect(slide, CONTENT_X, epr_y, CONTENT_W, 0.72,
             fill=RGBColor(0x04, 0x2a, 0x1e), border=C_TEAL)
    add_text(slide, "EPR MANDATE TAILWIND",
        CONTENT_X + 0.12, epr_y + 0.06, CONTENT_W - 0.24, 0.22,
        size=8, bold=True, color=C_TEAL)
    add_text(slide,
        "CPCB mandates 5% RRBO blending in FY25, scaling to 50% by FY31. "
        "Santosh Group II+ RRBO is purpose-built to fulfil this mandate.",
        CONTENT_X + 0.12, epr_y + 0.28, CONTENT_W - 0.24, 0.38,
        size=8.5, color=C_MUTED)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OUR PROMOTERS (Pattern C — 4-column grid)
# ═══════════════════════════════════════════════════════════════════════════════

PROMOTERS = [
    ("Lalit Bindal", "Founder &\nManaging Director", [
        "35+ years in oil & energy sector",
        "IOCL SSI Stockist since 2003",
        "6-district SERVO network, Western UP",
    ]),
    ("Abhinav Bindal", "Director —\nEngineering & Operations", [
        "Structural Engineer, Nabih Youseef Associates, LA",
        "Expertise in large-scale infrastructure design",
        "Driving digital & commercial strategy",
    ]),
    ("Pooja Bindal", "Director —\nAdministration", [
        "8 years independent business management",
        "Proprietor, Bindal Creation (garments & retail)",
        "Financial oversight & administrative leadership",
    ]),
    ("Robin Kumar", "Director —\nOperations", [
        "Based in Noida, Gautam Budh Nagar",
        "Operations & logistics management",
        "Supply chain for feedstock collection",
    ]),
]

def build_slide_04(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_logo(slide, logo)

    # Full-width header
    add_text(slide, "LEADERSHIP TEAM",
        0.35, 0.22, 10.0, 0.25, size=8, bold=True, color=C_TEAL)
    add_text(slide, "Experienced Leadership with Deep IOCL Roots",
        0.35, 0.50, 12.0, 0.55, size=22, bold=True, color=C_WHITE)

    # 4 columns
    col_w = (W - 0.7) / 4 - 0.1
    col_y0 = 1.22
    col_h = H - col_y0 - 0.35
    cx0 = 0.35

    for i, (name, role, bullets) in enumerate(PROMOTERS):
        cx = cx0 + i * (col_w + 0.1)
        add_rect(slide, cx, col_y0, col_w, col_h, fill=C_CARD, border=C_CARD_BORDER)

        # Circular headshot placeholder
        circle_d = 0.85
        circle_x = cx + (col_w - circle_d) / 2
        circle_y = col_y0 + 0.15
        c = slide.shapes.add_shape(9,
            Inches(circle_x), Inches(circle_y),
            Inches(circle_d), Inches(circle_d))
        c.fill.solid(); c.fill.fore_color.rgb = C_GREY
        c.line.color.rgb = C_TEAL; c.line.width = Pt(1.5)

        # Name
        add_text(slide, name,
            cx + 0.08, col_y0 + 1.1, col_w - 0.16, 0.32,
            size=10.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Role
        add_text(slide, role,
            cx + 0.08, col_y0 + 1.42, col_w - 0.16, 0.5,
            size=8.5, color=C_TEAL, align=PP_ALIGN.CENTER)
        # Teal divider
        add_rect(slide, cx + 0.25, col_y0 + 1.96, col_w - 0.5, 0.03, fill=C_TEAL)
        # Bullets
        for j, b in enumerate(bullets):
            add_text(slide, f"• {b}",
                cx + 0.1, col_y0 + 2.06 + j * 0.52, col_w - 0.2, 0.5,
                size=8.5, color=C_MUTED)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE OPPORTUNITY (Pattern A)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_05(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_pattern_a_image(slide, images.get("photo-1611273426858-450d8e3c9fce"))
    add_logo(slide, logo)
    add_slide_tag(slide, "MARKET OPPORTUNITY")
    add_headline(slide, "India's Used Oil Crisis —\nA ₹28,000 Cr Market", top=0.55, size=21)

    # 2×2 stat grid
    stats = [
        ("1.3 MMT", "Used lubricating oil generated in India annually"),
        ("<15%", "Formally recycled — rest burned or dumped"),
        ("257", "CPCB-registered re-refiners in India"),
        ("~0", "Group II+ RRBO producers in North India"),
    ]
    chip_w = (CONTENT_W - 0.1) / 2
    chip_h = 1.0
    gap = 0.1
    grid_y = 1.72
    for i, (val, lbl) in enumerate(stats):
        col = i % 2; row = i // 2
        x = CONTENT_X + col * (chip_w + gap)
        y = grid_y + row * (chip_h + gap)
        bdr = C_TEAL if i in (0, 3) else C_CARD_BORDER
        add_rect(slide, x, y, chip_w, chip_h, fill=C_CARD, border=bdr)
        add_text(slide, val,
            x + 0.12, y + 0.1, chip_w - 0.24, 0.42,
            size=22, bold=True,
            color=C_TEAL if i in (0, 3) else C_WHITE)
        add_text(slide, lbl,
            x + 0.12, y + 0.52, chip_w - 0.24, 0.42,
            size=8.5, color=C_MUTED)

    # Insight callout
    ins_y = grid_y + 2 * (chip_h + gap) + 0.12
    add_rect(slide, CONTENT_X, ins_y, CONTENT_W, 0.82,
             fill=RGBColor(0x04, 0x2a, 0x1e), border=C_TEAL)
    add_text(slide, "THE SANTOSH ADVANTAGE",
        CONTENT_X + 0.12, ins_y + 0.07, CONTENT_W - 0.24, 0.22,
        size=8, bold=True, color=C_TEAL)
    add_text(slide,
        "Only Group II+ producer in North India. Hydrotreating eliminates sulfur to "
        "<300 ppm — competitors (IFP) stuck at Group I via adsorption. EPR-compliant output.",
        CONTENT_X + 0.12, ins_y + 0.30, CONTENT_W - 0.24, 0.46,
        size=9, color=C_MUTED)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MARKET SIZE & GROWTH (Pattern A)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_06(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_pattern_a_image(slide, images.get("photo-1486325212027-8081e485255e"))
    add_logo(slide, logo)
    add_slide_tag(slide, "MARKET ANALYSIS")
    add_headline(slide, "High Growth,\nPolicy-Backed Market", top=0.55, size=24)

    # 4 KPI chips (2×2)
    kpis = [
        ("USD 3.38B", "India waste oil market (2025)"),
        ("~7.9% CAGR", "Growth to USD 6.22B by 2033"),
        ("4.2 MMT", "India lubricant market (FY22)"),
        ("5–6 MMT", "Projected lubricant market FY31"),
    ]
    chip_w = (CONTENT_W - 0.1) / 2
    chip_h = 0.85
    for i, (val, lbl) in enumerate(kpis):
        col = i % 2; row = i // 2
        x = CONTENT_X + col * (chip_w + 0.1)
        y = 1.68 + row * (chip_h + 0.1)
        add_kpi_chip(slide, val, lbl, x, y, width=chip_w, height=chip_h)

    # Market drivers
    drivers = [
        ("Regulatory Push", "EPR mandate drives formal collection; informal burning criminalised"),
        ("Rising Vehicle Ownership", "Auto sector growth = more used oil generated annually"),
        ("Import Substitution", "Group II+ RRBO replaces costly virgin base oil imports"),
        ("IOCL-Re MOU (Mar 2026)", "India's first national Group II+ RRBO circular economy framework"),
    ]
    driver_y = 1.68 + 2 * (chip_h + 0.1) + 0.15
    for i, (title, desc) in enumerate(drivers):
        add_bullet_row(slide, "•", title, desc,
                       top=driver_y + i * 0.56,
                       icon_color=C_TEAL if i < 2 else RGBColor(0x00, 0x6e, 0x56))

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PLANT & INFRASTRUCTURE (Pattern A)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_08(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_pattern_a_image(slide, images.get("photo-1518709268805-4e9042af9f23"))
    add_logo(slide, logo)
    add_slide_tag(slide, "INFRASTRUCTURE")
    add_headline(slide, "World-Class Plant —\nSCADA-Controlled, CPCB-Compliant", top=0.55, size=20)

    sections = [
        ("A. Distillation", "Dehydration column, vacuum distillation, WFE (REVA), asphalt collector"),
        ("B. Hydrotreating", "3 reactors with metallic catalysts, heat exchangers, H₂ recycle compressor"),
        ("C. Fractionation", "Fractionation column, pre-heaters, pressure vessels, pumps & coolers"),
        ("D. Tank Farm", "Storage: 1000 KL×2, 500 KL×2, 200 KL×2, 100 KL×2 + vapor recovery"),
        ("E. Utilities", "Thermal oil heater, cooling tower, compressed air, 4 chemical injection systems"),
        ("F. Control System", "SCADA + PLC plant-wide, remote monitoring, VFDs on all pumps"),
    ]
    row_h = 0.62
    row_y0 = 1.62
    for i, (sec, desc) in enumerate(sections):
        y = row_y0 + i * (row_h + 0.06)
        fill = RGBColor(0x12, 0x1e, 0x2c) if i % 2 == 0 else C_CARD
        add_rect(slide, CONTENT_X, y, CONTENT_W, row_h, fill=fill)
        add_rect(slide, CONTENT_X, y, 0.04, row_h, fill=C_TEAL)
        add_text(slide, sec,
            CONTENT_X + 0.14, y + 0.07, 1.65, 0.26,
            size=9, bold=True, color=C_TEAL)
        add_text(slide, desc,
            CONTENT_X + 0.14, y + 0.33, CONTENT_W - 0.24, 0.26,
            size=8.5, color=C_MUTED)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — LOCATION STRATEGY (Pattern A)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_09(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_pattern_a_image(slide, images.get("photo-1567427017947-545c5f8d16ad"))
    add_logo(slide, logo)
    add_slide_tag(slide, "LOCATION STRATEGY")
    add_headline(slide, "Ghaziabad — Heart of India's\nIndustrial Lubricant Belt", top=0.55, size=20)
    add_subheadline(slide,
        "Proximity to NCR, Western UP, and major lubricant blending hubs",
        top=1.42)

    points = [
        ("NCR Proximity",
         "Delhi-Meerut corridor — access to India's largest auto & industrial cluster"),
        ("IOCL Distribution Network",
         "Existing 6-district SERVO network = immediate feedstock & customer reach"),
        ("Western UP Industrial Belt",
         "Ghaziabad, Meerut, Muzaffarnagar, Baghpat — consistent used oil volumes"),
        ("Highway Connectivity",
         "NH-58 & NH-9 — excellent tanker logistics for inbound feedstock & outbound RRBO"),
        ("Used Oil Catchment",
         "30M+ vehicles in NCR alone; 200,000+ KL used oil/yr in catchment area"),
        ("Competitor Gap",
         "IFP (Sahibabad) is only re-refiner nearby — cannot produce Group II+"),
    ]
    pt_y = 1.90
    for title, desc in points:
        add_bullet_row(slide, "•", title, desc, top=pt_y)
        pt_y += 0.55

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PROJECT TIMELINE (Pattern B — Full-Width)
# ═══════════════════════════════════════════════════════════════════════════════

PHASES = [
    ("Pre-Project", "Year 1", "COMPLETE", [
        "Company incorporation & DPR",
        "Technology partner (REVA) selected",
    ], True),
    ("Construction", "Year 2", "COMPLETE", [
        "Civil works, foundations, buildings",
        "REVA plant equipment ordered & delivered",
    ], True),
    ("Commissioning", "Year 3 — NOW", "IN PROGRESS", [
        "Plant commissioning underway",
        "CPCB registration & QC lab (27 instruments)",
    ], False),
    ("Full Operations", "Year 4", "PLANNED", [
        "65 TPD full capacity",
        "Revenue target: ₹169 Cr+",
    ], False),
    ("Scale-Up", "Year 5–7", "PLANNED", [
        "Expand to 200 TPD capacity",
        "Pan-India RRBO supply capability",
    ], False),
]

def build_slide_10(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_logo(slide, logo)
    add_text(slide, "PROJECT STATUS",
        0.35, 0.22, 10.0, 0.25, size=8, bold=True, color=C_TEAL)
    add_text(slide, "Commissioning in Progress — Production Imminent",
        0.35, 0.50, 12.0, 0.55, size=22, bold=True, color=C_WHITE)

    phase_w = (W - 0.7) / 5 - 0.08
    phase_h = 4.6
    phase_y = 1.2
    px0 = 0.35

    for i, (name, period, status, bullets, done) in enumerate(PHASES):
        px = px0 + i * (phase_w + 0.08)
        active = (status == "IN PROGRESS")

        # Top status bar color
        bar_color = C_TEAL if done else (RGBColor(0x00, 0x6e, 0x56) if active else C_CARD_BORDER)
        add_rect(slide, px, phase_y, phase_w, 0.07, fill=bar_color)

        # Phase box
        box_fill = RGBColor(0x09, 0x2e, 0x22) if active else C_CARD
        add_rect(slide, px, phase_y + 0.07, phase_w, phase_h - 0.07,
                 fill=box_fill,
                 border=C_TEAL if active else C_CARD_BORDER)

        # Phase name
        add_text(slide, name,
            px + 0.1, phase_y + 0.15, phase_w - 0.2, 0.5,
            size=10, bold=True, color=C_WHITE)
        add_text(slide, period,
            px + 0.1, phase_y + 0.65, phase_w - 0.2, 0.25,
            size=8, color=C_MUTED)

        # Status badge
        badge_color = C_TEAL if done else (RGBColor(0xf3, 0x9c, 0x12) if active else C_CARD)
        badge_text_color = C_BG if done else (C_BG if active else C_MUTED)
        add_rect(slide, px + 0.1, phase_y + 0.95, phase_w - 0.2, 0.3,
                 fill=badge_color, border=None)
        add_text(slide, ("✓ " if done else "> " if active else "") + status,
            px + 0.12, phase_y + 0.97, phase_w - 0.24, 0.26,
            size=7.5, bold=True, color=badge_text_color, align=PP_ALIGN.CENTER)

        # Bullets
        for j, b in enumerate(bullets):
            add_text(slide, f"• {b}",
                px + 0.1, phase_y + 1.38 + j * 0.68, phase_w - 0.2, 0.65,
                size=8.5, color=C_MUTED)

        # Connector arrow
        if i < 4:
            add_text(slide, "▶",
                px + phase_w + 0.01, phase_y + phase_h / 2 - 0.1, 0.07, 0.25,
                size=8, color=C_TEAL, align=PP_ALIGN.CENTER)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — GROWTH SCALE-UP (Pattern B — Full-Width 3-Phase)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_13(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_logo(slide, logo)
    add_text(slide, "SCALE-UP PLAN",
        0.35, 0.22, 10.0, 0.25, size=8, bold=True, color=C_TEAL)
    add_text(slide, "65 TPD Today  →  200 TPD Tomorrow",
        0.35, 0.50, 12.0, 0.55, size=24, bold=True, color=C_WHITE)
    add_text(slide, "Designed for scale from day one — infrastructure ready for 3× expansion",
        0.35, 1.08, 11.0, 0.28, size=10, color=C_MUTED)

    phases = [
        ("PHASE 1", "Current", "COMMISSIONING", [
            ("Feed Capacity", "65 TPD"),
            ("RRBO Output", "~16,000 KL/yr"),
            ("Revenue Target", "₹169 Cr by Year 4"),
        ], True),
        ("PHASE 2", "Year 5–6", "PLANNED", [
            ("Feed Capacity", "130 TPD"),
            ("RRBO Output", "~32,000 KL/yr"),
            ("Revenue Target", "₹320–340 Cr"),
        ], False),
        ("PHASE 3", "Year 7+", "DESIGNED CAPACITY", [
            ("Feed Capacity", "200 TPD"),
            ("RRBO Output", "~50,000 KL/yr"),
            ("Revenue Target", "₹450–500 Cr"),
        ], False),
    ]
    col_w = (W - 0.7) / 3 - 0.12
    col_h = 4.5
    col_y = 1.44
    cx0 = 0.35

    for i, (phase, period, status, rows, active) in enumerate(phases):
        cx = cx0 + i * (col_w + 0.12)
        fill = RGBColor(0x09, 0x2e, 0x22) if active else C_CARD
        border = C_TEAL if active else C_CARD_BORDER
        top_fill = C_TEAL if active else RGBColor(0x00, 0x3d, 0x2e)
        add_rect(slide, cx, col_y, col_w, 0.07, fill=top_fill)
        add_rect(slide, cx, col_y + 0.07, col_w, col_h, fill=fill, border=border)
        add_text(slide, phase,
            cx + 0.15, col_y + 0.16, col_w - 0.3, 0.35,
            size=13, bold=True, color=C_TEAL if active else C_WHITE)
        add_text(slide, period,
            cx + 0.15, col_y + 0.52, col_w - 0.3, 0.25,
            size=9, color=C_MUTED)
        add_text(slide, status,
            cx + 0.15, col_y + 0.82, col_w - 0.3, 0.25,
            size=8, bold=True, color=C_TEAL if active else C_MUTED)
        for j, (lbl, val) in enumerate(rows):
            y = col_y + 1.22 + j * 1.0
            add_text(slide, lbl,
                cx + 0.15, y, col_w - 0.3, 0.26,
                size=8.5, color=C_MUTED)
            add_text(slide, val,
                cx + 0.15, y + 0.28, col_w - 0.3, 0.42,
                size=16, bold=True, color=C_TEAL if active else C_WHITE)
        if i < 2:
            add_text(slide, "→",
                cx + col_w + 0.02, col_y + col_h / 2 - 0.1, 0.1, 0.3,
                size=14, bold=True, color=C_TEAL)

    # EPR callout bar
    epr_y = col_y + col_h + 0.18
    add_rect(slide, 0.35, epr_y, W - 0.7, 0.52,
             fill=RGBColor(0x04, 0x1d, 0x13), border=C_CARD_BORDER)
    add_text(slide,
        "EPR mandate scaling from 5% (FY25) → 50% (FY31) will create demand for 650,000+ KL "
        "of RRBO annually — far exceeding current formal supply. Santosh is positioned to be "
        "North India's primary Group II+ supplier.",
        0.48, epr_y + 0.09, W - 1.0, 0.35,
        size=8.5, color=C_MUTED)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — PARTNERSHIP PROPOSAL / ASK (Pattern B — Full-Width 3-col)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_16(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_logo(slide, logo)
    add_text(slide, "IOCL PARTNERSHIP ASK",
        0.35, 0.22, 12.0, 0.25, size=8, bold=True, color=C_TEAL)
    add_text(slide, "Three Strategic Pillars of Partnership",
        0.35, 0.50, 11.5, 0.55, size=24, bold=True, color=C_WHITE)
    add_text(slide,
        "Together: India's first integrated Group II+ RRBO circular economy in North India",
        0.35, 1.08, 11.5, 0.28, size=10, color=C_MUTED)

    pillars = [
        ("01", "Offtake Agreement\n— RRBO Supply", [
            "Preferred supplier for Group II+ RRBO to IOCL blending plants",
            "3–5 year structured offtake, volume-tiered, aligned to EPR ramp",
            "SN 150, SN 500, Bright Stock grades available",
            "Enables IOCL SERVO to meet CPCB EPR mandates",
        ]),
        ("02", "Feedstock Collection\n— Used Oil Network", [
            "Leverage IOCL SERVO dealer network as used oil collection points",
            "IOCL dealers earn collection incentives; IOCL earns EPR credit",
            "Santosh provides containers, logistics, EPR certificates",
            "Formalised, traceable feedstock pipeline — replaces informal burning",
        ]),
        ("03", "Strategic Investment\n& Technical Support", [
            "IOCL equity or debt participation to accelerate 200 TPD scale-up",
            "Technical validation from IOCL R&D (Faridabad)",
            "Joint branding: IOCL-certified Group II+ RRBO",
            "LOI/MOU to unlock bank financing at preferential rates",
        ]),
    ]
    col_w = (W - 0.7) / 3 - 0.1
    col_h = 5.0
    col_y = 1.44
    cx0 = 0.35

    for i, (num, title, bullets) in enumerate(pillars):
        cx = cx0 + i * (col_w + 0.1)
        add_rect(slide, cx, col_y, col_w, 0.07, fill=C_TEAL)
        add_rect(slide, cx, col_y + 0.07, col_w, col_h, fill=C_CARD, border=C_CARD_BORDER)
        add_text(slide, num,
            cx + 0.15, col_y + 0.16, col_w - 0.3, 0.45,
            size=22, bold=True, color=C_TEAL)
        add_text(slide, title,
            cx + 0.15, col_y + 0.62, col_w - 0.3, 0.75,
            size=11, bold=True, color=C_WHITE)
        add_rect(slide, cx + 0.15, col_y + 1.44, col_w - 0.3, 0.03,
                 fill=C_CARD_BORDER)
        for j, b in enumerate(bullets):
            add_text(slide, f"• {b}",
                cx + 0.15, col_y + 1.55 + j * 0.78, col_w - 0.3, 0.72,
                size=9, color=C_MUTED)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CAPITAL INVESTMENT (Pattern A)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_11(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_pattern_a_image(slide, images.get("photo-1611532736597-de2d4265fba3"))
    add_logo(slide, logo)
    add_slide_tag(slide, "CAPITAL STRUCTURE")
    add_headline(slide, "Total Project Outlay:\n₹159.87 Crore", top=0.55, size=24)
    add_subheadline(slide, "30% Promoter Equity  |  70% Term Loan  |  9% Interest Rate", top=1.48)

    sources = [
        ("Promoters' Contribution", "30%", "₹46.62 Cr", "Self-funded equity"),
        ("Term Loan (Plant)", "70%", "₹94.79 Cr", "Bank financing for plant"),
        ("Working Capital Loan", "70% of WC", "₹14.00 Cr", "Operational working capital"),
    ]
    card_h = 1.1
    card_y0 = 2.00  # shifted from 1.90 to avoid overlap with subheadline (0.02" gap)
    for i, (title, pct, amount, note) in enumerate(sources):
        y = card_y0 + i * (card_h + 0.12)
        add_rect(slide, CONTENT_X, y, CONTENT_W, card_h, fill=C_CARD, border=C_CARD_BORDER)
        add_rect(slide, CONTENT_X, y, 0.04, card_h, fill=C_TEAL)
        add_text(slide, title,
            CONTENT_X + 0.14, y + 0.1, CONTENT_W * 0.55, 0.28,
            size=10, bold=True, color=C_WHITE)
        add_text(slide, pct,
            CONTENT_X + 0.14, y + 0.42, 1.0, 0.28,
            size=9, color=C_MUTED)
        add_text(slide, amount,
            CONTENT_X + CONTENT_W - 1.8, y + 0.1, 1.7, 0.42,
            size=18, bold=True, color=C_TEAL, align=PP_ALIGN.RIGHT)
        add_text(slide, note,
            CONTENT_X + 0.14, y + 0.72, CONTENT_W - 0.28, 0.28,
            size=8.5, color=C_MUTED)

    # Interest note
    add_text(slide, "Interest Rate: 9% p.a. (Term Loan & Working Capital)",
        CONTENT_X, card_y0 + 3 * (card_h + 0.12) + 0.1, CONTENT_W, 0.28,
        size=8.5, color=C_MUTED, italic=True)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — FINANCIAL PROJECTIONS (Pattern A — drawn bar chart left)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_12(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_logo(slide, logo)

    # ── Left panel: drawn bar chart ──────────────────────────────────────────
    chart_left = 0.35
    chart_w = IMG_W - 0.35
    chart_bottom = H - 0.55
    chart_h = 4.2
    chart_top = chart_bottom - chart_h

    years = ["Y3", "Y4", "Y5", "Y6", "Y7"]
    revenues = [122, 169, 200, 240, 280]  # ₹ Cr (approximate projections)
    max_rev = 300
    bar_gap = 0.14
    bar_w = (chart_w - (len(years) + 1) * bar_gap) / len(years)

    # Chart background
    add_rect(slide, chart_left, chart_top - 0.1, chart_w, chart_h + 0.2,
             fill=RGBColor(0x07, 0x12, 0x1c))
    add_text(slide, "REVENUE PROJECTION (₹ Crore)",
        chart_left + 0.1, chart_top - 0.08, chart_w - 0.2, 0.24,
        size=7.5, bold=True, color=C_MUTED)

    for i, (yr, rev) in enumerate(zip(years, revenues)):
        bx = chart_left + bar_gap + i * (bar_w + bar_gap)
        bar_h_actual = (rev / max_rev) * (chart_h - 0.5)
        by = chart_bottom - bar_h_actual - 0.25

        fill = C_TEAL if yr == "Y4" else RGBColor(0x00, 0x6e, 0x56)
        add_rect(slide, bx, by, bar_w, bar_h_actual, fill=fill)
        add_text(slide, f"₹{rev}Cr",
            bx, by - 0.3, bar_w, 0.28,
            size=7.5, bold=True,
            color=C_TEAL if yr == "Y4" else C_MUTED,
            align=PP_ALIGN.CENTER)
        add_text(slide, yr,
            bx, chart_bottom - 0.22, bar_w, 0.22,
            size=8, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Right-edge fade
    add_rect(slide, IMG_W - 0.5, 0, 0.5, H, fill=C_BG)

    # ── Right panel ──────────────────────────────────────────────────────────
    add_slide_tag(slide, "FINANCIAL MODEL")
    add_headline(slide, "Revenue ₹169 Cr+\nby Year 4", top=0.55, size=24)

    prod_data = [
        ("Year 3", "46 KL/day", "15,243 KL/yr RRBO"),
        ("Year 4+", "65 KL/day", "21,775 KL/yr RRBO"),
    ]
    add_text(slide, "PRODUCTION PROFILE",
        CONTENT_X, 1.68, CONTENT_W, 0.24,
        size=7.5, bold=True, color=C_MUTED)
    for i, (yr, daily, annual) in enumerate(prod_data):
        add_kpi_chip(slide, daily, f"{yr} — {annual}",
                     CONTENT_X + i * (CONTENT_W / 2 + 0.02), 1.96,
                     width=CONTENT_W / 2 - 0.05, height=0.9)

    pricing = [
        ("RRBO Pricing", "₹80,000/KL (Y3)  →  ₹88,305/KL (Y7)"),
        ("Feedstock Cost", "₹45,000/KL (Y3, rising ~3%/yr)"),
        ("Gross Margin", "~43–44% on RRBO sales"),
    ]
    ratio_y = 3.06
    add_text(slide, "KEY METRICS",
        CONTENT_X, ratio_y, CONTENT_W, 0.24,
        size=7.5, bold=True, color=C_MUTED)
    for i, (lbl, val) in enumerate(pricing):
        y = ratio_y + 0.28 + i * 0.6
        add_rect(slide, CONTENT_X, y, CONTENT_W, 0.55,
                 fill=C_CARD if i % 2 == 0 else RGBColor(0x12, 0x1e, 0x2c))
        add_text(slide, lbl,
            CONTENT_X + 0.1, y + 0.08, 1.55, 0.22,
            size=8.5, bold=True, color=C_TEAL)
        add_text(slide, val,
            CONTENT_X + 1.7, y + 0.08, CONTENT_W - 1.8, 0.40,
            size=8.5, color=C_WHITE)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — FEEDSTOCK PIPELINE (Pattern A)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_14(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_pattern_a_image(slide, images.get("photo-1519003722824-194d4455a60c"))
    add_logo(slide, logo)
    add_slide_tag(slide, "FEEDSTOCK STRATEGY")
    add_headline(slide, "Building a Reliable\nFeedstock Network", top=0.55, size=23)

    sources = [
        ("Automotive Workshops", "Car/truck service centres — largest used engine oil source in NCR"),
        ("Industrial Plants", "Manufacturing units using hydraulic, gear and cutting oils"),
        ("Transport Fleets", "Commercial trucking & bus depots — regular oil change cycles"),
        ("IOCL Dealer Network", "SERVO stockist relationships enable formal collection"),
        ("EPR-Registered Collectors", "CPCB-registered aggregators incentivised to deliver to re-refiners"),
    ]
    src_y = 1.68
    for title, desc in sources:
        add_bullet_row(slide, "•", title, desc, top=src_y)
        src_y += 0.56

    # 4 KPI chips
    kpis = [
        ("15,243 KL", "Year 3 Feed Target"),
        ("21,775 KL", "Year 4+ Feed Target"),
        ("₹45,000/KL", "Feedstock Cost (Y3)"),
        ("150–200 km", "Collection Radius"),
    ]
    chip_w = (CONTENT_W - 0.09) / 2
    chip_h = 0.72
    kpi_y = src_y + 0.1
    for i, (val, lbl) in enumerate(kpis):
        col = i % 2; row = i // 2
        add_kpi_chip(slide,
            val, lbl,
            CONTENT_X + col * (chip_w + 0.09), kpi_y + row * (chip_h + 0.09),
            width=chip_w, height=chip_h)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — STRATEGIC PARTNERSHIPS (Pattern A)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_15(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_pattern_a_image(slide, images.get("photo-1450101499163-c8848c66ca85"))
    add_logo(slide, logo)
    add_slide_tag(slide, "PARTNERSHIPS")
    add_headline(slide, "Built on Trusted\nRelationships", top=0.55, size=25)

    partners = [
        ("REVA Process Technologies, Pune", "Technology Partner", [
            "European vacuum distillation + hydrotreating technology",
            "70–75% base oil recovery guarantee; ongoing catalyst supply",
        ]),
        ("Indian Oil Corporation Ltd.", "23+ Year Partner (SSI since 2003) + IOCL-Re MOU Mar 2026", [
            "6-district SERVO distribution network; deep trust established",
            "MOU creates framework for Group II+ RRBO circular economy",
        ]),
        ("Hindustan Petroleum Corp. Ltd.", "O&M & Bottling Partner", [
            "LPG Bottling: Amroha — 30 TMT/yr; O&M: Sitarganj — 340 TMT tenure",
            "Validates Santosh's capacity to operate large-scale energy assets",
        ]),
    ]
    card_h = 1.45
    card_y0 = 1.72
    for i, (name, role, bullets) in enumerate(partners):
        y = card_y0 + i * (card_h + 0.1)
        add_rect(slide, CONTENT_X, y, CONTENT_W, card_h, fill=C_CARD, border=C_CARD_BORDER)
        add_rect(slide, CONTENT_X, y, 0.04, card_h, fill=C_TEAL)
        add_text(slide, name,
            CONTENT_X + 0.14, y + 0.1, CONTENT_W - 0.24, 0.28,
            size=10, bold=True, color=C_TEAL)
        add_text(slide, role,
            CONTENT_X + 0.14, y + 0.38, CONTENT_W - 0.24, 0.26,
            size=8.5, color=C_MUTED, italic=True)
        for j, b in enumerate(bullets):
            add_text(slide, f"• {b}",
                CONTENT_X + 0.14, y + 0.66 + j * 0.36, CONTENT_W - 0.28, 0.34,
                size=8.5, color=C_WHITE)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — VALUE PROP FOR IOCL (Pattern A)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_17(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_pattern_a_image(slide, images.get("photo-1558618666-fcd25c85cd64"))
    add_logo(slide, logo)
    add_slide_tag(slide, "IOCL VALUE PROP")
    add_headline(slide, "Why This Partnership\nDelivers for IOCL", top=0.55, size=23)

    values = [
        ("EPR Compliance — Zero CapEx",
         "SERVO must blend RRBO per CPCB mandate. Source from Santosh, skip building own re-refinery."),
        ("Captive Group II+ Supply in North India",
         "No certified Group II+ RRBO producer in North India. Santosh fills this gap."),
        ("Circular Economy Credential",
         "IOCL-Re MOU (Mar 2026) commits to Group II+ circular economy. Santosh makes it real."),
        ("SERVO Dealer Network Monetisation",
         "Dealers become collection points, earn per-litre fees. Deepens dealer relationships."),
        ("Import Substitution Alignment",
         "Group II+ currently imported. Domestic RRBO from Santosh reduces forex outflow."),
        ("Trusted Promoter — 23+ Year Track Record",
         "Lalit Bindal has been an IOCL SSI for 23+ years. This is an existing partnership."),
    ]
    card_h = 0.75
    card_y0 = 1.70
    gap = 0.08
    for i, (title, desc) in enumerate(values):
        row = i % 3; col = i // 3
        x = CONTENT_X + col * (CONTENT_W / 2 + gap / 2)
        y = card_y0 + row * (card_h + gap)
        add_rect(slide, x, y, CONTENT_W / 2 - gap / 2, card_h,
                 fill=C_CARD, border=C_CARD_BORDER)
        add_rect(slide, x, y, 0.03, card_h, fill=C_TEAL)
        add_text(slide, title,
            x + 0.1, y + 0.06, CONTENT_W / 2 - 0.2, 0.25,
            size=8.5, bold=True, color=C_WHITE)
        add_text(slide, desc,
            x + 0.1, y + 0.34, CONTENT_W / 2 - 0.2, 0.36,
            size=7.5, color=C_MUTED)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — CTA / NEXT STEPS (Pattern A)
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_18(prs, logo, images):
    slide = new_slide(prs)
    add_bg(slide)
    add_pattern_a_image(slide, images.get("photo-1473341304170-971dccb5ac1e"))
    add_logo(slide, logo)
    add_slide_tag(slide, "NEXT STEPS")
    add_headline(slide,
        "Let's Build India's\nCircular Oil Economy Together",
        top=0.45, size=21)

    steps = [
        ("1", "Site Visit",
         "IOCL team visits Ghaziabad plant — inspect commissioning progress & technology"),
        ("2", "Technical Review",
         "Share DPR, lab data, REVA process docs with IOCL R&D (Faridabad)"),
        ("3", "MOU / Letter of Intent",
         "Non-binding MOU covering RRBO offtake volumes, feedstock, investment framework"),
        ("4", "Commercial Agreement",
         "Binding offtake agreement, pricing formula, EPR transfer mechanism — target Q2 FY26"),
    ]
    step_h = 0.85
    step_y0 = 1.65  # shifted from 1.55 to avoid touching headline bottom (0" gap)
    for i, (num, title, desc) in enumerate(steps):
        y = step_y0 + i * (step_h + 0.1)
        add_rect(slide, CONTENT_X, y, CONTENT_W, step_h, fill=C_CARD, border=C_CARD_BORDER)
        # Number badge
        badge = slide.shapes.add_shape(9,
            Inches(CONTENT_X + 0.1), Inches(y + 0.17),
            Inches(0.48), Inches(0.48))
        badge.fill.solid(); badge.fill.fore_color.rgb = C_TEAL
        badge.line.fill.background()
        add_text(slide, num,
            CONTENT_X + 0.1, y + 0.17, 0.48, 0.48,
            size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        add_text(slide, title,
            CONTENT_X + 0.72, y + 0.1, CONTENT_W - 0.82, 0.28,
            size=10, bold=True, color=C_WHITE)
        add_text(slide, desc,
            CONTENT_X + 0.72, y + 0.42, CONTENT_W - 0.82, 0.38,
            size=8.5, color=C_MUTED)

    # Contact block
    contact_y = step_y0 + 4 * (step_h + 0.1) + 0.1
    add_rect(slide, CONTENT_X, contact_y, CONTENT_W, 0.85,
             fill=RGBColor(0x04, 0x2a, 0x1e), border=C_TEAL)
    add_text(slide, "Lalit Bindal — Managing Director",
        CONTENT_X + 0.15, contact_y + 0.08, CONTENT_W - 0.3, 0.26,
        size=10, bold=True, color=C_WHITE)
    add_text(slide,
        "santoshgzb@yahoo.com   ·   +91 98101 21438   ·   www.santosh-petrochemical.com",
        CONTENT_X + 0.15, contact_y + 0.38, CONTENT_W - 0.3, 0.38,
        size=8.5, color=C_MUTED)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

# All Unsplash photo IDs used in the deck
IMAGE_IDS = [
    "photo-1565008576549-57569a49371d",  # slide 01: refinery aerial
    "photo-1504328345606-18bbc8c9d7d1",  # slide 02: LPG tanker/ops
    "photo-1518709268805-4e9042af9f23",  # slides 03, 08: industrial plant
    "photo-1611273426858-450d8e3c9fce",  # slide 05: oil drums
    "photo-1486325212027-8081e485255e",  # slide 06: industrial growth
    "photo-1567427017947-545c5f8d16ad",  # slide 09: NCR highway/aerial
    "photo-1611532736597-de2d4265fba3",  # slide 11: finance/industry
    "photo-1519003722824-194d4455a60c",  # slide 14: oil tanker truck
    "photo-1450101499163-c8848c66ca85",  # slide 15: handshake
    "photo-1558618666-fcd25c85cd64",     # slide 17: lubricant/blending
    "photo-1473341304170-971dccb5ac1e",  # slide 18: circular economy
]

SLIDE_BUILDERS = [
    build_slide_01, build_slide_02, build_slide_03, build_slide_04,
    build_slide_05, build_slide_06, build_slide_07, build_slide_08,
    build_slide_09, build_slide_10, build_slide_11, build_slide_12,
    build_slide_13, build_slide_14, build_slide_15, build_slide_16,
    build_slide_17, build_slide_18,
]

OUTPUT_PATH = "asset/Santosh_IOCL_Pitch_Deck_v2.pptx"


def verify_output(path):
    """Sanity-check the generated deck."""
    from pptx import Presentation as _Prs
    prs = _Prs(path)
    assert len(prs.slides) == 18, f"Expected 18 slides, got {len(prs.slides)}"
    assert abs(prs.slide_width.inches - 13.3) < 0.1
    assert abs(prs.slide_height.inches - 7.5) < 0.1
    for i, slide in enumerate(prs.slides):
        assert len(slide.shapes) >= 3, f"Slide {i+1} has too few shapes ({len(slide.shapes)})"
    print("✅ 18 slides verified")
    print("✅ Dimensions 13.3\" × 7.5\" verified")
    print("✅ All slides have content shapes")


def main():
    print("=== IOCL Pitch Deck Rebuilder ===\n")

    # 1. Convert logo
    print("Converting logo...")
    logo = get_logo_bytes()

    # 2. Download images
    print("\nDownloading images...")
    images = {}
    for pid in IMAGE_IDS:
        images[pid] = download_image(pid)

    # 3. Build presentation
    print("\nBuilding slides...")
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    for i, builder in enumerate(SLIDE_BUILDERS):
        print(f"  Slide {i+1:02d}: {builder.__name__}")
        builder(prs, logo, images)

    # 4. Save
    prs.save(OUTPUT_PATH)
    print(f"\nSaved: {OUTPUT_PATH}")

    # 5. Verify
    print("\nVerifying...")
    verify_output(OUTPUT_PATH)
    print(f"\n✅ Done! Open with: open \"{OUTPUT_PATH}\"")


if __name__ == "__main__":
    main()
