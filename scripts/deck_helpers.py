# scripts/deck_helpers.py
"""Design system constants and shape-builder helpers for the IOCL pitch deck."""

import subprocess, os, requests
from io import BytesIO
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colors ────────────────────────────────────────────────────────────────────
C_BG          = RGBColor(0x0f, 0x19, 0x23)  # slide background
C_TEAL        = RGBColor(0x00, 0xb8, 0x94)  # accent primary
C_TEAL2       = RGBColor(0x00, 0xce, 0xc9)  # accent secondary
C_WHITE       = RGBColor(0xff, 0xff, 0xff)  # text primary
C_MUTED       = RGBColor(0x8a, 0x8f, 0x98)  # text muted
C_CARD        = RGBColor(0x14, 0x22, 0x32)  # card fill
C_CARD_BORDER = RGBColor(0x00, 0x5a, 0x4a)  # card border
C_FOOTER      = RGBColor(0x07, 0x0f, 0x18)  # footer bar
C_GREY        = RGBColor(0x44, 0x44, 0x55)  # placeholder circles

FONT = "Calibri"

# ── Slide constants ────────────────────────────────────────────────────────────
W  = 13.3   # slide width (inches)
H  = 7.5    # slide height (inches)
IMG_W = 5.59          # Pattern A image panel width (42% of W)
CONTENT_X = 5.85      # Pattern A content panel left edge
CONTENT_W = W - CONTENT_X - 0.2   # content panel width
FOOTER_Y  = H - 0.22  # footer top
FOOTER_H  = 0.22
LOGO_X    = W - 0.6   # logo left
LOGO_Y    = 0.08
LOGO_SIZE = 0.45


# ── Shape builders ─────────────────────────────────────────────────────────────

def add_bg(slide):
    """Full-slide dark background rectangle."""
    sh = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(W), Inches(H))
    sh.fill.solid(); sh.fill.fore_color.rgb = C_BG
    sh.line.fill.background()
    return sh


def add_rect(slide, left, top, width, height, fill=None, border=None, border_pt=0.5):
    """Add a solid rectangle. fill/border are RGBColor or None (transparent)."""
    sh = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(border_pt)
    else:
        sh.line.fill.background()
    return sh


def add_text(slide, text, left, top, width, height,
             size=11, bold=False, color=None, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    """Add a styled text box."""
    if color is None:
        color = C_WHITE
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT
    return tb


def add_text_multiline(slide, lines, left, top, width, height,
                        size=10, bold=False, color=None, spacing_pt=4):
    """Add multiple lines as separate paragraphs in one text box."""
    if color is None:
        color = C_MUTED
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing_pt if i > 0 else 0)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return tb


def add_image(slide, img_bytes_or_path, left, top, width, height):
    """Embed an image from bytes or file path."""
    if isinstance(img_bytes_or_path, (bytes, bytearray)):
        stream = BytesIO(img_bytes_or_path)
    else:
        stream = img_bytes_or_path  # already a BytesIO
    slide.shapes.add_picture(stream, Inches(left), Inches(top), Inches(width), Inches(height))


def add_footer(slide):
    """Standard footer bar with company name + teal dot."""
    add_rect(slide, 0, FOOTER_Y, W, FOOTER_H, fill=C_FOOTER)
    add_text(slide,
        "Santosh Petrochemical Innovations Pvt. Ltd.  ·  Confidential",
        0.15, FOOTER_Y + 0.02, 9.0, FOOTER_H - 0.04,
        size=6.5, color=C_MUTED)
    dot = slide.shapes.add_shape(9,
        Inches(W - 0.25), Inches(FOOTER_Y + 0.05),
        Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = C_TEAL
    dot.line.fill.background()


def add_logo(slide, logo_bytes):
    """Place logo PNG top-right. Falls back to teal circle with 'S'."""
    if logo_bytes:
        stream = BytesIO(logo_bytes)
        slide.shapes.add_picture(stream,
            Inches(LOGO_X), Inches(LOGO_Y),
            Inches(LOGO_SIZE), Inches(LOGO_SIZE))
    else:
        c = slide.shapes.add_shape(9,
            Inches(LOGO_X), Inches(LOGO_Y),
            Inches(LOGO_SIZE), Inches(LOGO_SIZE))
        c.fill.solid(); c.fill.fore_color.rgb = C_TEAL
        c.line.fill.background()
        add_text(slide, "S",
            LOGO_X, LOGO_Y, LOGO_SIZE, LOGO_SIZE,
            size=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)


def add_slide_tag(slide, text):
    """Teal ALL-CAPS tag on right panel."""
    add_text(slide, text.upper(),
        CONTENT_X, 0.28, CONTENT_W, 0.22,
        size=8, color=C_TEAL, bold=True)


def add_headline(slide, text, top=0.55, size=26):
    """Large bold white headline on right panel."""
    add_text(slide, text,
        CONTENT_X, top, CONTENT_W, 1.1,
        size=size, bold=True, color=C_WHITE)


def add_subheadline(slide, text, top=1.35):
    """Smaller muted subheadline."""
    add_text(slide, text,
        CONTENT_X, top, CONTENT_W, 0.4,
        size=11, color=C_MUTED)


def add_kpi_chip(slide, value, label, left, top, width=1.55, height=0.75):
    """KPI card: teal value + muted label."""
    add_rect(slide, left, top, width, height, fill=C_CARD, border=C_CARD_BORDER)
    add_text(slide, value,
        left + 0.08, top + 0.08, width - 0.16, 0.38,
        size=17, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, label,
        left + 0.06, top + 0.44, width - 0.12, 0.26,
        size=8, color=C_MUTED, align=PP_ALIGN.CENTER)


def add_bullet_row(slide, icon_char, title, desc, top, icon_color=None):
    """One bullet row: colored title + muted description."""
    if icon_color is None:
        icon_color = C_TEAL
    add_rect(slide, CONTENT_X, top, 0.04, 0.45, fill=icon_color)
    add_text(slide, title,
        CONTENT_X + 0.12, top, CONTENT_W - 0.12, 0.22,
        size=10, bold=True, color=C_WHITE)
    add_text(slide, desc,
        CONTENT_X + 0.12, top + 0.22, CONTENT_W - 0.12, 0.24,
        size=9, color=C_MUTED)


def add_card(slide, left, top, width, height, title, lines,
             title_size=10, line_size=9, border_color=None):
    """Info card with title + lines of text."""
    if border_color is None:
        border_color = C_CARD_BORDER
    add_rect(slide, left, top, width, height, fill=C_CARD, border=border_color)
    add_text(slide, title,
        left + 0.1, top + 0.08, width - 0.2, 0.25,
        size=title_size, bold=True, color=C_TEAL)
    for i, line in enumerate(lines):
        add_text(slide, line,
            left + 0.1, top + 0.32 + i * 0.28, width - 0.2, 0.26,
            size=line_size, color=C_MUTED)


def add_pattern_a_image(slide, img_bytes):
    """Left image panel for Pattern A slides."""
    if img_bytes:
        add_image(slide, img_bytes, 0, 0, IMG_W, H)
    else:
        add_rect(slide, 0, 0, IMG_W, H, fill=RGBColor(0x06, 0x2a, 0x22))
    add_rect(slide, IMG_W - 0.5, 0, 0.5, H, fill=C_BG)


# ── Image utilities ────────────────────────────────────────────────────────────

CACHE_DIR = "/tmp/deck_images"
os.makedirs(CACHE_DIR, exist_ok=True)

def download_image(photo_id):
    """Download Unsplash image by photo ID. Returns bytes or None on failure."""
    cache_path = os.path.join(CACHE_DIR, f"{photo_id}.jpg")
    if os.path.exists(cache_path):
        with open(cache_path, "rb") as f:
            return f.read()
    url = f"https://images.unsplash.com/{photo_id}?w=1200&q=85"
    try:
        resp = requests.get(url, timeout=15)
        resp.raise_for_status()
        data = resp.content
        with open(cache_path, "wb") as f:
            f.write(data)
        print(f"  ✓ Downloaded {photo_id}")
        return data
    except Exception as e:
        print(f"  ⚠ Failed to download {photo_id}: {e}")
        return None


def get_logo_bytes():
    """Convert SVG logo to PNG bytes using rsvg-convert."""
    svg_path = "asset/Santosh Petrochemical logo.svg"
    png_path = "/tmp/deck_logo.png"
    try:
        subprocess.run(
            ["/opt/homebrew/bin/rsvg-convert",
             "-w", "200", "-h", "200",
             svg_path, "-o", png_path],
            check=True, capture_output=True
        )
        with open(png_path, "rb") as f:
            print("  ✓ Logo converted")
            return f.read()
    except Exception as e:
        print(f"  ⚠ Logo conversion failed: {e}. Using fallback.")
        return None
